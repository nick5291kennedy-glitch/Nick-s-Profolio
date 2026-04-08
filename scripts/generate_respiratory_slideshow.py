from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "assets" / "respiratory_images"
GENERATED_DIR = ASSET_DIR / "generated"
OUTPUT = ROOT / "respiratory_system_nursing_slideshow.pptx"


PALETTE = {
    "navy": "10324A",
    "teal": "1E6F7A",
    "aqua": "5CC8C6",
    "sky": "D9F2F0",
    "cream": "F6FBFB",
    "gold": "E8B04B",
    "coral": "D86C61",
    "ink": "18313E",
    "mist": "E9F4F3",
    "white": "FFFFFF",
}

TITLE_FONT = "Avenir Next"
BODY_FONT = "Avenir Next"


def rgb(hex_code: str) -> RGBColor:
    return RGBColor.from_string(hex_code)


def ensure_dirs() -> None:
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)


def font(size: int, bold: bool = False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Avenir Next.ttc",
        "/System/Library/Fonts/SFNS.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size=size, index=1 if bold else 0)
            except Exception:
                try:
                    return ImageFont.truetype(path, size=size)
                except Exception:
                    continue
    return ImageFont.load_default()


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font_obj, width: int):
    lines = []
    current = []
    for word in text.split():
        trial = " ".join(current + [word])
        if draw.textlength(trial, font=font_obj) <= width:
            current.append(word)
        else:
            if current:
                lines.append(" ".join(current))
            current = [word]
    if current:
        lines.append(" ".join(current))
    return lines


def create_canvas():
    return Image.new("RGB", (1600, 900), f"#{PALETTE['cream']}")


def soft_background(img):
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 1600, 900), fill=f"#{PALETTE['cream']}")
    draw.ellipse((1100, -80, 1650, 420), fill=f"#{PALETTE['sky']}")
    draw.ellipse((-180, 540, 360, 1080), fill=f"#{PALETTE['mist']}")
    draw.rounded_rectangle((70, 70, 1530, 830), radius=44, outline=f"#{PALETTE['sky']}", width=4)
    return draw


def save_image(name: str, painter):
    img = create_canvas()
    draw = soft_background(img)
    painter(img, draw)
    out = GENERATED_DIR / name
    img.save(out)
    return out


def make_ventilation_image():
    def painter(img, draw):
        title_f = font(52, bold=True)
        sub_f = font(28)
        draw.text((110, 90), "Ventilation Cycle", font=title_f, fill=f"#{PALETTE['navy']}")
        draw.text((112, 160), "Diaphragm movement changes thoracic volume and pressure.", font=sub_f, fill=f"#{PALETTE['teal']}")

        panels = [(120, 260, 730, 760), (870, 260, 1480, 760)]
        labels = [
            ("Inspiration", "Diaphragm contracts and flattens\nThoracic cavity expands\nAir pressure falls, air moves in"),
            ("Expiration", "Diaphragm relaxes and rises\nElastic recoil decreases volume\nAir pressure rises, air moves out"),
        ]
        for idx, box in enumerate(panels):
            x1, y1, x2, y2 = box
            draw.rounded_rectangle(box, radius=36, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
            cx = (x1 + x2) // 2
            draw.ellipse((cx - 115, y1 + 95, cx + 115, y1 + 310), outline=f"#{PALETTE['teal']}", width=10)
            draw.line((cx, y1 + 250, cx, y1 + 420), fill=f"#{PALETTE['teal']}", width=12)
            if idx == 0:
                draw.arc((cx - 180, y1 + 260, cx + 180, y1 + 520), 200, 340, fill=f"#{PALETTE['coral']}", width=12)
                draw.polygon([(cx - 220, y1 + 185), (cx - 280, y1 + 165), (cx - 280, y1 + 205)], fill=f"#{PALETTE['gold']}")
                draw.line((cx - 160, y1 + 185, cx - 260, y1 + 185), fill=f"#{PALETTE['gold']}", width=10)
            else:
                draw.arc((cx - 180, y1 + 180, cx + 180, y1 + 440), 20, 160, fill=f"#{PALETTE['coral']}", width=12)
                draw.polygon([(cx + 220, y1 + 185), (cx + 280, y1 + 165), (cx + 280, y1 + 205)], fill=f"#{PALETTE['gold']}")
                draw.line((cx + 160, y1 + 185, cx + 260, y1 + 185), fill=f"#{PALETTE['gold']}", width=10)

            draw.text((x1 + 34, y1 + 24), labels[idx][0], font=font(38, bold=True), fill=f"#{PALETTE['navy']}")
            bullet_y = y1 + 540
            for line in labels[idx][1].split("\n"):
                draw.text((x1 + 42, bullet_y), f"• {line}", font=font(26), fill=f"#{PALETTE['ink']}")
                bullet_y += 44

    return save_image("ventilation_cycle.png", painter)


def make_gas_exchange_image():
    def painter(img, draw):
        draw.text((100, 85), "Alveolar Gas Exchange", font=font(52, bold=True), fill=f"#{PALETTE['navy']}")
        draw.text((102, 155), "Diffusion follows pressure gradients across a thin membrane.", font=font(28), fill=f"#{PALETTE['teal']}")
        draw.rounded_rectangle((100, 250, 760, 760), radius=42, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
        draw.rounded_rectangle((840, 250, 1500, 760), radius=42, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)

        draw.ellipse((180, 320, 690, 720), fill=f"#{PALETTE['mist']}", outline=f"#{PALETTE['aqua']}", width=8)
        draw.ellipse((250, 390, 620, 650), outline=f"#{PALETTE['teal']}", width=8)
        draw.text((285, 500), "ALVEOLUS", font=font(44, bold=True), fill=f"#{PALETTE['navy']}")

        for y in range(350, 690, 70):
            draw.line((820, y, 1490, y), fill=f"#{PALETTE['sky']}", width=10)
        for x in range(920, 1450, 120):
            draw.ellipse((x, 395, x + 95, 470), fill=f"#F5A09A", outline=f"#{PALETTE['coral']}", width=4)
            draw.text((x + 26, 414), "RBC", font=font(18, bold=True), fill=f"#{PALETTE['white']}")

        draw.line((675, 470, 920, 470), fill=f"#{PALETTE['gold']}", width=12)
        draw.polygon([(920, 470), (885, 445), (885, 495)], fill=f"#{PALETTE['gold']}")
        draw.text((745, 410), "O₂ to blood", font=font(34, bold=True), fill=f"#{PALETTE['gold']}")

        draw.line((920, 590, 680, 590), fill=f"#{PALETTE['teal']}", width=12)
        draw.polygon([(680, 590), (715, 565), (715, 615)], fill=f"#{PALETTE['teal']}")
        draw.text((760, 625), "CO₂ to alveoli", font=font(34, bold=True), fill=f"#{PALETTE['teal']}")

        draw.rounded_rectangle((820, 770, 1490, 860), radius=28, fill=f"#{PALETTE['navy']}")
        draw.text((850, 790), "Nursing focus: watch for hypoxemia, fatigue, and worsening work of breathing.", font=font(24), fill=f"#{PALETTE['white']}")

    return save_image("gas_exchange_flow.png", painter)


def make_vq_image():
    def painter(img, draw):
        draw.text((100, 85), "V/Q Balance and Control of Breathing", font=font(48, bold=True), fill=f"#{PALETTE['navy']}")
        draw.text((102, 150), "Matching airflow with blood flow is the key to oxygenation.", font=font(28), fill=f"#{PALETTE['teal']}")
        draw.rounded_rectangle((110, 240, 760, 760), radius=40, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
        draw.rounded_rectangle((840, 240, 1490, 760), radius=40, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)

        draw.text((150, 280), "Ventilation-Perfusion", font=font(36, bold=True), fill=f"#{PALETTE['navy']}")
        draw.rectangle((180, 390, 240, 620), fill=f"#{PALETTE['aqua']}")
        draw.rectangle((310, 330, 370, 680), fill=f"#{PALETTE['teal']}")
        draw.rectangle((440, 430, 500, 580), fill=f"#{PALETTE['gold']}")
        draw.text((168, 640), "Shunt", font=font(24, bold=True), fill=f"#{PALETTE['ink']}")
        draw.text((260, 700), "Ideal match", font=font(24, bold=True), fill=f"#{PALETTE['ink']}")
        draw.text((425, 600), "Dead\nspace", font=font(24, bold=True), fill=f"#{PALETTE['ink']}")
        draw.text((560, 420), "Low V, high Q:\npneumonia, atelectasis", font=font(24), fill=f"#{PALETTE['ink']}")
        draw.text((560, 555), "High V, low Q:\npulmonary embolism", font=font(24), fill=f"#{PALETTE['ink']}")

        draw.text((880, 280), "Neural and Chemical Control", font=font(36, bold=True), fill=f"#{PALETTE['navy']}")
        nodes = [
            ((950, 395, 1135, 500), "Medulla\nsets rhythm", PALETTE["teal"]),
            ((1180, 395, 1365, 500), "Pons\nmodulates depth", PALETTE["aqua"]),
            ((950, 565, 1135, 670), "Central\nCO₂ / pH", PALETTE["gold"]),
            ((1180, 565, 1365, 670), "Peripheral\nPaO₂ / PaCO₂ / pH", PALETTE["coral"]),
        ]
        for box, label, color in nodes:
            draw.rounded_rectangle(box, radius=22, fill=f"#{color}", outline=f"#{PALETTE['white']}", width=3)
            draw.multiline_text((box[0] + 28, box[1] + 24), label, font=font(28, bold=True), fill=f"#{PALETTE['white']}", spacing=4)
        draw.line((1135, 448, 1180, 448), fill=f"#{PALETTE['navy']}", width=6)
        draw.line((1042, 500, 1042, 565), fill=f"#{PALETTE['navy']}", width=6)
        draw.line((1272, 500, 1272, 565), fill=f"#{PALETTE['navy']}", width=6)

    return save_image("vq_balance.png", painter)


def make_assessment_image():
    def painter(img, draw):
        draw.text((100, 85), "Respiratory Assessment Flow", font=font(50, bold=True), fill=f"#{PALETTE['navy']}")
        draw.text((102, 150), "A systematic assessment helps nurses identify early deterioration.", font=font(28), fill=f"#{PALETTE['teal']}")
        steps = [
            ("1", "Inspect", "Rate, pattern, accessory muscles,\ncyanosis, mental status"),
            ("2", "Palpate", "Chest expansion, tracheal position,\nsubcutaneous emphysema"),
            ("3", "Percuss", "Resonance, dullness,\nhyperresonance"),
            ("4", "Auscultate", "Vesicular sounds,\ncrackles, wheezes, rhonchi"),
        ]
        x = 120
        for num, title, desc in steps:
            draw.rounded_rectangle((x, 290, x + 300, 710), radius=28, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
            draw.ellipse((x + 30, 325, x + 110, 405), fill=f"#{PALETTE['teal']}")
            draw.text((x + 56, 347), num, font=font(34, bold=True), fill=f"#{PALETTE['white']}")
            draw.text((x + 30, 445), title, font=font(34, bold=True), fill=f"#{PALETTE['navy']}")
            lines = desc.split("\n")
            y = 510
            for line in lines:
                draw.text((x + 30, y), line, font=font(24), fill=f"#{PALETTE['ink']}")
                y += 38
            if x < 1080:
                draw.line((x + 300, 500, x + 340, 500), fill=f"#{PALETTE['gold']}", width=8)
                draw.polygon([(x + 340, 500), (x + 314, 480), (x + 314, 520)], fill=f"#{PALETTE['gold']}")
            x += 340

    return save_image("assessment_flow.png", painter)


def make_abg_image():
    def painter(img, draw):
        draw.text((100, 85), "Pulse Oximetry and ABG Quick Look", font=font(48, bold=True), fill=f"#{PALETTE['navy']}")
        draw.text((102, 150), "Use SpO₂ trends together with ABGs and the full clinical picture.", font=font(28), fill=f"#{PALETTE['teal']}")
        draw.rounded_rectangle((110, 250, 610, 780), radius=38, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
        draw.text((150, 290), "ABG Normals", font=font(38, bold=True), fill=f"#{PALETTE['navy']}")
        rows = [("pH", "7.35 - 7.45"), ("PaCO₂", "35 - 45 mmHg"), ("HCO₃⁻", "22 - 26 mEq/L"), ("PaO₂", "80 - 100 mmHg"), ("SaO₂", "95 - 100%")]
        y = 360
        for label, value in rows:
            draw.rounded_rectangle((145, y, 575, y + 72), radius=20, fill=f"#{PALETTE['mist']}")
            draw.text((175, y + 18), label, font=font(30, bold=True), fill=f"#{PALETTE['navy']}")
            draw.text((360, y + 18), value, font=font(28), fill=f"#{PALETTE['ink']}")
            y += 88

        draw.rounded_rectangle((690, 250, 1490, 780), radius=38, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
        draw.text((735, 290), "Interpretation Tips", font=font(38, bold=True), fill=f"#{PALETTE['navy']}")
        tips = [
            ("SpO₂", "Reflects oxygen saturation, not CO₂ removal."),
            ("Acidosis", "Think hypoventilation, fatigue, CNS depression."),
            ("Alkalosis", "Think anxiety, pain, sepsis, early hypoxemia."),
            ("Trend", "Compare symptoms, work of breathing, and response to care."),
        ]
        y = 380
        for label, desc in tips:
            draw.ellipse((735, y + 4, 770, y + 39), fill=f"#{PALETTE['gold']}")
            draw.text((790, y), label, font=font(28, bold=True), fill=f"#{PALETTE['teal']}")
            lines = wrap_text(draw, desc, font(24), 600)
            yy = y + 36
            for line in lines:
                draw.text((790, yy), line, font=font(24), fill=f"#{PALETTE['ink']}")
                yy += 30
            y = yy + 24

    return save_image("abg_quicklook.png", painter)


def make_devices_image():
    def painter(img, draw):
        draw.text((100, 85), "Oxygen Delivery Spectrum", font=font(50, bold=True), fill=f"#{PALETTE['navy']}")
        draw.text((102, 150), "Select the device that matches oxygen needs and monitoring intensity.", font=font(28), fill=f"#{PALETTE['teal']}")
        devices = [
            ("Nasal cannula", "1 - 6 L/min\nLow-flow"),
            ("Simple mask", "5 - 10 L/min\nShort-term"),
            ("Venturi mask", "Precise FiO₂\nCOPD support"),
            ("Nonrebreather", "High FiO₂\nAcute hypoxemia"),
            ("HFNC / Vent", "Advanced support\nEscalate quickly"),
        ]
        x = 95
        colors = [PALETTE["teal"], PALETTE["aqua"], PALETTE["gold"], PALETTE["coral"], PALETTE["navy"]]
        for idx, (title, desc) in enumerate(devices):
            draw.rounded_rectangle((x, 300, x + 270, 700), radius=28, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
            draw.rounded_rectangle((x + 32, 330, x + 238, 405), radius=20, fill=f"#{colors[idx]}")
            draw.text((x + 48, 350), title, font=font(26, bold=True), fill=f"#{PALETTE['white']}")
            draw.ellipse((x + 82, 455, x + 188, 560), outline=f"#{colors[idx]}", width=8)
            draw.line((x + 135, 560, x + 135, 620), fill=f"#{colors[idx]}", width=8)
            draw.multiline_text((x + 52, 615), desc, font=font(22), fill=f"#{PALETTE['ink']}", spacing=6, align="center")
            x += 295

    return save_image("oxygen_devices.png", painter)


def make_case_image():
    def painter(img, draw):
        draw.text((100, 85), "Clinical Snapshot", font=font(52, bold=True), fill=f"#{PALETTE['navy']}")
        draw.text((102, 150), "Link assessment data to immediate nursing priorities.", font=font(28), fill=f"#{PALETTE['teal']}")
        draw.rounded_rectangle((115, 250, 720, 780), radius=38, fill=f"#{PALETTE['white']}", outline=f"#{PALETTE['sky']}", width=4)
        draw.text((160, 300), "Patient: pneumonia with acute decline", font=font(34, bold=True), fill=f"#{PALETTE['navy']}")
        metrics = [("RR", "30/min"), ("SpO₂", "88% RA"), ("Temp", "38.7°C"), ("LOC", "Confused")]
        y = 390
        for label, value in metrics:
            draw.rounded_rectangle((160, y, 675, y + 78), radius=18, fill=f"#{PALETTE['mist']}")
            draw.text((190, y + 18), label, font=font(28, bold=True), fill=f"#{PALETTE['teal']}")
            draw.text((410, y + 18), value, font=font(28, bold=True), fill=f"#{PALETTE['ink']}")
            y += 98

        draw.rounded_rectangle((800, 250, 1490, 780), radius=38, fill=f"#{PALETTE['navy']}")
        draw.text((850, 300), "First Nursing Actions", font=font(36, bold=True), fill=f"#{PALETTE['white']}")
        actions = [
            "Raise head of bed and start supplemental oxygen.",
            "Complete focused respiratory assessment and vitals.",
            "Notify provider / rapid response for worsening status.",
            "Prepare for ABGs, cultures, imaging, and IV therapy.",
        ]
        y = 380
        for action in actions:
            draw.ellipse((850, y + 8, 878, y + 36), fill=f"#{PALETTE['gold']}")
            lines = wrap_text(draw, action, font(26), 530)
            yy = y
            for line in lines:
                draw.text((900, yy), line, font=font(26), fill=f"#{PALETTE['white']}")
                yy += 34
            y = yy + 28

    return save_image("clinical_snapshot.png", painter)


def build_graphics():
    return {
        "ventilation": make_ventilation_image(),
        "exchange": make_gas_exchange_image(),
        "vq": make_vq_image(),
        "assessment": make_assessment_image(),
        "abg": make_abg_image(),
        "devices": make_devices_image(),
        "case": make_case_image(),
    }


def set_slide_bg(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color_hex)


def add_textbox(slide, left, top, width, height, text, size=20, color="18313E", bold=False, font_name=BODY_FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_bullets(slide, left, top, width, height, bullets, size=21, color="18313E", accent="1E6F7A", line_space=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.line_spacing = line_space
        p.space_after = Pt(8)
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None, dark=False):
    title_color = "FFFFFF" if dark else PALETTE["navy"]
    subtitle_color = "D9F2F0" if dark else PALETTE["teal"]
    add_textbox(slide, Inches(0.7), Inches(0.45), Inches(8.8), Inches(0.9), title, size=28, color=title_color, bold=True, font_name=TITLE_FONT)
    if subtitle:
        add_textbox(slide, Inches(0.72), Inches(1.1), Inches(8.8), Inches(0.55), subtitle, size=13, color=subtitle_color, font_name=BODY_FONT)


def add_card(slide, left, top, width, height, fill_hex="FFFFFF", line_hex="D9F2F0"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.color.rgb = rgb(line_hex)
    shape.line.width = Pt(1.5)
    return shape


def add_image(slide, path, left, top, width=None, height=None):
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def build_presentation():
    ensure_dirs()
    graphics = build_graphics()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    lung_img = ASSET_DIR / "lung_anatomy.jpg"
    alveoli_img = ASSET_DIR / "alveoli.png"
    oxygen_img = ASSET_DIR / "oxygen_therapy.jpg"

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["navy"])
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.7), Inches(-0.7), Inches(5.2), Inches(3.4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(PALETTE["teal"])
    accent.line.fill.background()
    panel = add_card(slide, Inches(0.55), Inches(0.6), Inches(7.2), Inches(6.1), fill_hex="FFFFFF", line_hex=PALETTE["navy"])
    panel.fill.transparency = 0.05
    add_textbox(slide, Inches(0.95), Inches(1.0), Inches(5.8), Inches(1.0), "Respiratory System", size=30, color=PALETTE["navy"], bold=True, font_name=TITLE_FONT)
    add_textbox(slide, Inches(0.98), Inches(1.95), Inches(5.4), Inches(1.2), "College Nursing Level\nStructure, physiology, assessment, and priority care", size=20, color=PALETTE["teal"], bold=False)
    add_textbox(slide, Inches(0.98), Inches(4.95), Inches(4.7), Inches(0.5), "Prepared for classroom presentation and study review", size=12, color=PALETTE["ink"])
    stat = add_card(slide, Inches(0.98), Inches(5.55), Inches(3.05), Inches(0.7), fill_hex=PALETTE["mist"], line_hex=PALETTE["sky"])
    add_textbox(slide, Inches(1.18), Inches(5.75), Inches(2.7), Inches(0.24), "Focus: oxygenation, ventilation, and clinical nursing judgment", size=11, color=PALETTE["navy"], bold=True)
    add_image(slide, lung_img, Inches(8.2), Inches(1.05), width=Inches(4.2))

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Learning Objectives", "What nursing students should understand by the end of the slideshow")
    add_card(slide, Inches(0.65), Inches(1.7), Inches(6.1), Inches(4.9), fill_hex="FFFFFF")
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.0),
        Inches(5.5),
        Inches(4.1),
        [
            "Identify the major structures of the upper and lower respiratory tract.",
            "Explain how ventilation, perfusion, and diffusion support gas exchange.",
            "Connect pulse oximetry, ABGs, and physical assessment to patient status.",
            "Apply nursing priorities to common respiratory disorders and acute decline.",
        ],
        size=22,
    )
    add_image(slide, graphics["assessment"], Inches(7.0), Inches(1.75), width=Inches(5.5))

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Core Anatomy", "Key structures that support filtration, conduction, and gas exchange")
    add_image(slide, lung_img, Inches(0.75), Inches(1.65), width=Inches(4.2))
    add_card(slide, Inches(5.2), Inches(1.7), Inches(7.3), Inches(4.9), fill_hex="FFFFFF")
    add_bullets(
        slide,
        Inches(5.5),
        Inches(2.0),
        Inches(6.6),
        Inches(4.0),
        [
            "Upper airway includes the nose, sinuses, pharynx, and larynx, which warm and humidify air.",
            "Lower airway includes the trachea, bronchi, bronchioles, and alveoli.",
            "The right lung has 3 lobes; the left lung has 2 lobes and a cardiac notch.",
            "Pleura and negative intrathoracic pressure help keep the lungs expanded.",
        ],
        size=20,
    )
    cue = add_card(slide, Inches(5.55), Inches(5.55), Inches(3.65), Inches(0.78), fill_hex=PALETTE["navy"], line_hex=PALETTE["navy"])
    add_textbox(slide, Inches(5.78), Inches(5.79), Inches(3.2), Inches(0.25), "Clinical cue: alveoli are the primary gas-exchange units.", size=12, color=PALETTE["white"])

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Ventilation Mechanics", "How pressure changes move air in and out of the lungs")
    add_image(slide, graphics["ventilation"], Inches(0.7), Inches(1.55), width=Inches(12.0))

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Gas Exchange and Transport", "Oxygenation depends on alveoli, perfusion, and hemoglobin delivery")
    add_image(slide, alveoli_img, Inches(0.85), Inches(1.8), width=Inches(4.2))
    add_image(slide, graphics["exchange"], Inches(5.15), Inches(1.65), width=Inches(7.1))

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "V/Q Match and Control of Breathing", "Major concepts behind hypoxemia and respiratory drive")
    add_image(slide, graphics["vq"], Inches(0.7), Inches(1.6), width=Inches(12.0))

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Respiratory Assessment", "Bedside data that often signal trouble before lab values return")
    add_image(slide, graphics["assessment"], Inches(0.72), Inches(1.7), width=Inches(8.2))
    add_card(slide, Inches(9.25), Inches(1.82), Inches(3.3), Inches(4.9), fill_hex="FFFFFF")
    add_textbox(slide, Inches(9.55), Inches(2.1), Inches(2.6), Inches(0.4), "Red Flags", size=24, color=PALETTE["navy"], bold=True)
    add_bullets(
        slide,
        Inches(9.42),
        Inches(2.65),
        Inches(2.8),
        Inches(3.6),
        [
            "Increasing accessory muscle use",
            "Silent chest or diminished air entry",
            "Confusion, agitation, or lethargy",
            "Cyanosis or rapidly falling SpO₂",
        ],
        size=18,
    )

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Pulse Oximetry and ABGs", "Use trends and patient appearance, not just single numbers")
    add_image(slide, graphics["abg"], Inches(0.68), Inches(1.72), width=Inches(12.0))

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Common Respiratory Disorders", "High-yield conditions frequently encountered in nursing care")
    disorders = [
        ("Asthma", "Reversible airway inflammation and bronchoconstriction with wheeze and cough."),
        ("COPD", "Chronic airflow limitation, hyperinflation, and frequent exacerbations."),
        ("Pneumonia", "Infection with inflammation, consolidation, and impaired oxygen diffusion."),
        ("ARDS", "Severe inflammatory lung injury with refractory hypoxemia and low compliance."),
    ]
    positions = [(0.8, 1.8), (3.95, 1.8), (7.1, 1.8), (10.25, 1.8)]
    colors = [PALETTE["teal"], PALETTE["aqua"], PALETTE["gold"], PALETTE["coral"]]
    for (x, y), (title, desc), color_hex in zip(positions, disorders, colors):
        add_card(slide, Inches(x), Inches(y), Inches(2.7), Inches(3.85), fill_hex="FFFFFF")
        pill = add_card(slide, Inches(x + 0.18), Inches(y + 0.2), Inches(1.75), Inches(0.55), fill_hex=color_hex, line_hex=color_hex)
        add_textbox(slide, Inches(x + 0.34), Inches(y + 0.35), Inches(1.4), Inches(0.2), title, size=14, color=PALETTE["white"], bold=True)
        add_textbox(slide, Inches(x + 0.2), Inches(y + 1.0), Inches(2.25), Inches(2.3), desc, size=17, color=PALETTE["ink"])
    add_textbox(slide, Inches(0.92), Inches(6.1), Inches(11.4), Inches(0.35), "Nursing priority across disorders: recognize worsening oxygenation early and escalate promptly.", size=16, color=PALETTE["navy"], bold=True, align=PP_ALIGN.CENTER)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Oxygen Therapy and Priority Interventions", "Match the treatment device to the severity of respiratory compromise")
    add_image(slide, oxygen_img, Inches(0.8), Inches(1.85), width=Inches(2.6))
    add_image(slide, graphics["devices"], Inches(3.7), Inches(1.65), width=Inches(8.6))
    add_card(slide, Inches(0.8), Inches(4.85), Inches(2.6), Inches(1.3), fill_hex=PALETTE["navy"], line_hex=PALETTE["navy"])
    add_textbox(slide, Inches(1.0), Inches(5.1), Inches(2.2), Inches(0.7), "Safety: lowest effective FiO₂, humidify when needed, and follow oxygen fire precautions.", size=13, color=PALETTE["white"])

    # Slide 11
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["cream"])
    add_title(slide, "Clinical Application", "From patient cues to action: what the nurse does first")
    add_image(slide, graphics["case"], Inches(0.7), Inches(1.68), width=Inches(12.0))

    # Slide 12
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, PALETTE["navy"])
    add_textbox(slide, Inches(0.9), Inches(0.85), Inches(5.8), Inches(0.8), "Key Takeaways", size=28, color=PALETTE["white"], bold=True, font_name=TITLE_FONT)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(1.7),
        Inches(5.8),
        Inches(3.4),
        [
            "Effective respiration requires a patent airway, adequate ventilation, intact diffusion, and perfusion.",
            "Respiratory assessment findings often reveal deterioration before tests confirm it.",
            "Pulse oximetry, ABGs, and mental-status changes should always be interpreted together.",
            "Fast nursing interventions can prevent respiratory failure and improve outcomes.",
        ],
        size=21,
        color=PALETTE["white"],
    )
    add_card(slide, Inches(7.1), Inches(1.1), Inches(5.25), Inches(4.9), fill_hex=PALETTE["white"], line_hex=PALETTE["white"])
    add_textbox(slide, Inches(7.45), Inches(1.45), Inches(4.6), Inches(0.4), "References", size=24, color=PALETTE["navy"], bold=True)
    add_bullets(
        slide,
        Inches(7.35),
        Inches(2.0),
        Inches(4.6),
        Inches(3.4),
        [
            "Hinkle & Cheever. Brunner & Suddarth's Textbook of Medical-Surgical Nursing.",
            "Lewis et al. Medical-Surgical Nursing: Assessment and Management of Clinical Problems.",
            "American Association for Respiratory Care clinical resources.",
            "National Heart, Lung, and Blood Institute education materials.",
        ],
        size=16,
    )
    add_textbox(slide, Inches(0.95), Inches(6.55), Inches(11.2), Inches(0.3), "Designed for a college nursing audience with visual review aids for speaking or studying.", size=12, color=PALETTE["sky"], align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
